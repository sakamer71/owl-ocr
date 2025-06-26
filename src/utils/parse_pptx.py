#!/usr/bin/env python3
"""
extract_pptx.py

Extracts:
 - all slide text frames
 - all tables (output as HTML)
 - all text within embedded images via OCR
   (including WMF/EMF by converting to PNG via LibreOffice headless)

Requirements:
    pip install unstructured[all-docs] python-pptx
    sudo apt update && sudo apt install libreoffice-writer libreoffice-core
"""

import subprocess
import argparse
import os
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from unstructured.partition.pptx import partition_pptx
from unstructured.partition.image import partition_image
from unstructured.documents.elements import Table


def extract_pptx_text_tables_images(pptx_path, images_dir=None):
    """
    Returns:
        text_runs   : list of strings (all text from text boxes, titles, etc.)
        tables_html : list of HTML strings (one per table)
    """
    # 1) Partition PPTX into text + tables
    elements = partition_pptx(filename=pptx_path)
    text_runs = []
    tables_html = []

    for el in elements:
        if isinstance(el, Table):
            # Render table as HTML
            try:
                # Try to get HTML from metadata first
                html = getattr(el.metadata, 'text_as_html', None)
                if not html:
                    # Fallback to converting table data to HTML
                    html = "<table>"
                    if hasattr(el, 'metadata') and hasattr(el.metadata, 'text_as_html'):
                        html = el.metadata.text_as_html
                    else:
                        # Create basic HTML table from table data
                        html = "<table><tr><td>Table content</td></tr></table>"
            except Exception as e:
                print(f"Warning: failed to render table as HTML: {e}")
                html = "<table><tr><td>Table content (rendering failed)</td></tr></table>"
            tables_html.append(html)
        else:
            text_runs.append(el.text)

    # 2) Extract and OCR images (with WMF/EMF conversion)
    prs = Presentation(pptx_path)
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    # Access image data safely using getattr
                    img = getattr(shape, 'image', None)
                    if img is not None:
                        ext = img.ext.lower()
                        if images_dir is not None:
                            os.makedirs(images_dir, exist_ok=True)
                            img_filename = os.path.join(images_dir, f"slide{slide_idx}_img{shape_idx}.{ext}")
                        else:
                            img_filename = f"slide{slide_idx}_img{shape_idx}.{ext}"
                        with open(img_filename, "wb") as fp:
                            fp.write(img.blob)

                        # Convert WMF/EMF via LibreOffice
                        if ext in ("wmf", "emf"):
                            base = os.path.splitext(img_filename)[0]
                            png_filename = f"{base}.png"
                            try:
                                # Run LibreOffice conversion
                                result = subprocess.run([
                                    "libreoffice", "--headless",
                                    "--convert-to", "png",
                                    "--outdir", images_dir if images_dir else os.getcwd(),
                                    img_filename
                                ], check=True, capture_output=True, text=True)
                                
                                # Check if the PNG file was actually created
                                if os.path.isfile(png_filename):
                                    img_to_ocr = png_filename
                                    print(f"Successfully converted {img_filename} to {png_filename}")
                                else:
                                    print(f"Warning: LibreOffice conversion completed but {png_filename} not found")
                                    print(f"LibreOffice output: {result.stdout}")
                                    if result.stderr:
                                        print(f"LibreOffice errors: {result.stderr}")
                                    print(f"Skipping OCR for {img_filename} - conversion failed")
                                    # Skip OCR for this image since conversion failed
                                    continue
                            except subprocess.CalledProcessError as e:
                                print(f"Warning: LibreOffice conversion failed for {img_filename}: {e}")
                                print(f"LibreOffice output: {e.stdout}")
                                if e.stderr:
                                    print(f"LibreOffice errors: {e.stderr}")
                                print(f"Skipping OCR for {img_filename} - conversion failed")
                                continue
                            except Exception as e:
                                print(f"Warning: unexpected error during conversion of {img_filename}: {e}")
                                print(f"Skipping OCR for {img_filename} - conversion failed")
                                continue
                        else:
                            img_to_ocr = img_filename

                        # OCR the raster image
                        try:
                            img_elements = partition_image(filename=img_to_ocr)
                            for img_el in img_elements:
                                text_runs.append(img_el.text)
                        except Exception as e:
                            print(f"Warning: OCR failed for {img_to_ocr}: {e}")
                except Exception as e:
                    print(f"Warning: failed to extract image from shape: {e}")

    return text_runs, tables_html


def main(args=None):
    if args is None:
        # Called directly, parse command line arguments
        parser = argparse.ArgumentParser(
            description="Extract text, tables, and OCR (including WMF/EMF) from a PPTX file"
        )
        parser.add_argument(
            "-i", "--input", dest="pptx_file", required=True,
            help="Path to input .pptx file"
        )
        parser.add_argument(
            "-o", "--out_dir", default="parsed_docs",
            help="Directory to write output files (default: parsed_docs)"
        )
        parser.add_argument(
            "-v", "--verbose", action="store_true",
            help="Enable verbose output"
        )
        parser.add_argument(
            "--stdout", action="store_true",
            help="Output text content to stdout instead of writing to files"
        )
        args = parser.parse_args()

    # Set verbosity based on args
    verbose = getattr(args, 'verbose', False)
    stdout_mode = getattr(args, 'stdout', False)
    
    # Process the PPTX file
    source_path = Path(args.pptx_file)
    base_name = source_path.stem  # filename without extension

    # In stdout mode, we still need a temp dir for images during processing
    # but we won't write the text files to disk
    if not stdout_mode:
        # Ensure output directory exists
        out_dir = Path(args.out_dir)
        out_dir.mkdir(exist_ok=True)
        
        out_text = out_dir / f"{base_name}.txt"
        out_tables = out_dir / f"{base_name}_tables.html"
        images_dir = out_dir / base_name
        
        if verbose:
            print(f"Processing PPTX file: {args.pptx_file}")
            print(f"Output directory: {out_dir}")
    else:
        # For stdout mode, we still need a temp dir for images
        import tempfile
        temp_dir = tempfile.mkdtemp(prefix="ocr_convert_")
        images_dir = Path(temp_dir) / base_name
        
        if verbose:
            print(f"Processing PPTX file: {args.pptx_file}")
            print(f"Temp directory for images: {images_dir}")
    
    texts, tables = extract_pptx_text_tables_images(args.pptx_file, images_dir=str(images_dir))

    if stdout_mode:
        # Output text content to stdout
        for t in texts:
            print(t.strip())
            print()
            
        if tables and verbose:
            print("\n=== Tables (HTML format) ===")
            for html in tables:
                print(html)
                print()
    else:
        # Write out all text runs to file
        with open(out_text, "w", encoding="utf-8") as f:
            for t in texts:
                f.write(t.strip() + "\n\n")

        # Write out tables as HTML
        with open(out_tables, "w", encoding="utf-8") as f:
            for html in tables:
                f.write(html + "\n\n")

        print(f"Done. Text → {out_text}; Tables → {out_tables}; Images → {images_dir}")


if __name__ == "__main__":
    main()

